import os
import subprocess
import shutil
import uuid
import textwrap
from flask import Flask, request, jsonify, render_template, session, send_from_directory
from flask_session import Session
from transformers import pipeline
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from gtts import gTTS
from PIL import Image, ImageDraw, ImageFont

app = Flask(__name__) 

# Load summarization pipeline
summarizer = pipeline('summarization', model='facebook/bart-large-cnn', device=0)

def summarize_text(text):
    """Summarize the input text."""
    summary = summarizer(text, max_length=150, min_length=100, do_sample=False)
    return summary[0]['summary_text']

def create_pptx(topic, summary, code, pptx_path):
    """Create a PowerPoint presentation."""
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "Made By RIYA CHAND"

    # Summary Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"
    content = slide.placeholders[1]
    content.text = summary
    format_text(content.text_frame, font_size=Pt(18))

    # Code Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Code"
    content = slide.placeholders[1]
    content.text = code
    format_text(content.text_frame, font_size=Pt(16))

    prs.save(pptx_path)

def format_text(text_frame, font_size=Pt(18)):
    """Format text in a text frame."""
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = font_size
        paragraph.alignment = PP_ALIGN.LEFT

def convert_pptx_to_images(pptx_path, output_dir):
    """Convert PPTX slides to images."""
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        img_path = os.path.join(output_dir, f'slide_{i+1}.png')
        img = Image.new('RGB', (1280, 720), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        title = slide.shapes.title.text if slide.shapes.title else ""
        content = slide.placeholders[1].text if slide.placeholders[1] else ""

        try:
            font = ImageFont.truetype("arial.ttf", 24)
        except IOError:
            font = ImageFont.load_default()

        draw_text(draw, title, (10, 10), font, (0, 0, 0))
        draw_text(draw, content, (10, 50), font, (0, 0, 0))
        img.save(img_path)

def draw_text(draw, text, position, font, color):
    """Draw wrapped text on an image."""
    wrapped_text = textwrap.fill(text, width=40)
    draw.text(position, wrapped_text, font=font, fill=color)

def create_audio_for_summary(summary, output_dir):
    """Create an audio file for the summary."""
    audio_path = os.path.join(output_dir, 'summary.mp3')
    tts = gTTS(text=summary, lang='en')
    tts.save(audio_path)
    return audio_path

def create_audio_for_slides(pptx_path, output_dir):
    """Create audio for each slide's notes."""
    prs = Presentation(pptx_path)
    audio_files = []
    for i, slide in enumerate(prs.slides):
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                audio_path = os.path.join(output_dir, f'slide_{i+1}.mp3')
                tts = gTTS(text=notes, lang='en')
                tts.save(audio_path)
                audio_files.append(audio_path)
    return audio_files

def create_video_from_images_and_audio(image_dir, audio_files, video_path):
    """Create a video from images and audio."""
    temp_video = os.path.join(image_dir, 'temp_video.mp4')
    
    # Create video from images
    image_cmd = [
        'ffmpeg', '-y', '-framerate', '1/5',  # 1 frame every 5 seconds
        '-pattern_type', 'glob',
        '-i', os.path.join(image_dir, '*.png'),
        '-s', '1920x1080',  # Full HD resolution
        '-c:v', 'libx264', '-pix_fmt', 'yuv420p', temp_video
    ]
    subprocess.run(image_cmd, check=True)

    concat_list_path = os.path.join(image_dir, 'concat_list.txt')
    with open(concat_list_path, 'w') as f:
        for audio_file in audio_files:
            video_segment = os.path.join(image_dir, f'video_{os.path.basename(audio_file)}.mp4')
            audio_cmd = [
                'ffmpeg', '-y', '-i', temp_video, '-i', audio_file,
                '-c:v', 'copy', '-c:a', 'aac', '-shortest', video_segment
            ]
            subprocess.run(audio_cmd, check=True)
            f.write(f"file '{video_segment}'\n")

    concat_cmd = ['ffmpeg', '-y', '-f', 'concat', '-safe', '0', '-i', concat_list_path, '-c', 'copy', video_path]
    subprocess.run(concat_cmd, check=True)

    return True

@app.route('/templates/<filename>')
def download_file(filename):
    """Serve the generated video for download."""
    temp_dir = os.path.join(os.getcwd(), 'temp')  # Directory where files are stored
    return send_from_directory(temp_dir, filename, as_attachment=True)

@app.route('/')
def index():
    """Render the homepage."""
    return render_template('index.html')

@app.route('/generate-video', methods=['POST'])
def generate_video():
    """Generate video from input data."""

    data = request.json
    topic = data.get('topic')
    content = data.get('content')
    code = data.get('code')

    if not topic or not content or not code:
        return jsonify({"error": "Topic, content, and code are required"}), 400

    summary = summarize_text(content)
    pptx_path = os.path.join(os.getcwd(), 'temp', 'presentation.pptx')
    image_dir = os.path.join(os.getcwd(), 'temp', 'images')
    video_path = os.path.join(os.getcwd(), 'temp', 'final_video.mp4')
    
    create_pptx(topic, summary, code, pptx_path)
    convert_pptx_to_images(pptx_path, image_dir)

    summary_audio = create_audio_for_summary(summary, image_dir)
    slide_audio_files = create_audio_for_slides(pptx_path, image_dir)
    audio_files = [summary_audio] + slide_audio_files

    if not create_video_from_images_and_audio(image_dir, audio_files, video_path):
        return jsonify({"error": "Video creation failed"}), 500

    video_filename = 'final_video.mp4'
    return jsonify({"message": "Video created successfully", "video_file": os.path.basename(video_path)})

@app.route('/cleanup', methods=['POST'])
def cleanup():
    """Cleanup session files."""
    # Session cleanup if needed
    return jsonify({"message": "Session cleaned up."})

if __name__ == '__main__':  
    app.run(debug=True)
